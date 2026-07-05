import os
import sys
import qrcode
import openpyxl
from pptx import Presentation
import win32com.client
import shutil

# Konfigurasi path
BASE_DIR = r"C:\sertifikat_digital"
EXCEL_PATH = os.path.join(BASE_DIR, "sertifikat.xlsx")
PPTX_TEMPLATE_PATH = os.path.join(BASE_DIR, "sertifikat.pptx")
OUTPUT_DIR = os.path.join(BASE_DIR, "output")
PDF_DIR = os.path.join(OUTPUT_DIR, "pdf")

# Github Pages base URL
GITHUB_BASE_URL = "https://shajara-ai.github.io/sertifikat"

def sanitize_filename(name):
    """Mengubah karakter tidak valid untuk nama file menjadi underscore"""
    invalid_chars = ['<', '>', ':', '"', '/', '\\', '|', '?', '*']
    sanitized = name
    for char in invalid_chars:
        sanitized = sanitized.replace(char, '_')
    return sanitized

def generate_qr_code(data, output_path):
    """Menghasilkan QR Code berupa gambar PNG"""
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_M,
        box_size=10,
        border=2,
    )
    qr.add_data(data)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    img.save(output_path)

def replace_text_in_shape(shape, replacements):
    """Mengganti teks placeholder di dalam shape PowerPoint dengan tetap mempertahankan format font"""
    if not shape.has_text_frame:
        return
    
    for paragraph in shape.text_frame.paragraphs:
        # Gabungkan teks paragraph untuk mengecek kecocokan placeholder
        text = paragraph.text
        has_key = False
        for key in replacements:
            if key in text:
                has_key = True
                break
        
        if has_key:
            # Lakukan penggantian teks
            for key, val in replacements.items():
                text = text.replace(key, str(val))
            
            # Update teks run pertama dan bersihkan run sisanya agar format font tidak rusak/hilang
            if paragraph.runs:
                paragraph.runs[0].text = text
                for run in paragraph.runs[1:]:
                    run.text = ""
            else:
                paragraph.text = text

def process_mail_merge():
    print("=== Memulai Proses Mail Merge Sertifikat ===")
    
    # Pastikan direktori output dan pdf ada
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    os.makedirs(PDF_DIR, exist_ok=True)
    
    if not os.path.exists(EXCEL_PATH):
        print(f"Error: File Excel tidak ditemukan di {EXCEL_PATH}")
        return
    if not os.path.exists(PPTX_TEMPLATE_PATH):
        print(f"Error: Template PowerPoint tidak ditemukan di {PPTX_TEMPLATE_PATH}")
        return

    # Muat database Excel
    try:
        wb = openpyxl.load_workbook(EXCEL_PATH)
        ws = wb.active
    except PermissionError:
        print("\n[ERROR] GAGAL MEMBACA EXCEL!")
        print("File 'sertifikat.xlsx' sedang dibuka di program lain (misal: Microsoft Excel).")
        print("Silakan tutup program Excel terlebih dahulu, kemudian jalankan kembali skrip ini.")
        return

    # Ambil header kolom untuk mapping indeks dari Baris 1
    headers = [cell.value for cell in ws[1]]
    print("Header terdeteksi:", [h for h in headers if h is not None])

    # Temukan indeks untuk setiap kolom dari Excel asli pengguna
    try:
        idx_timestamp = headers.index("Timestamp")
        idx_id = headers.index("Nomor Sertifikat")
        idx_penerima = headers.index("Penerima")
        idx_kegiatan = headers.index("Kegiatan")
        idx_dokumen = headers.index("Dokumen Link")
        idx_sebagai = headers.index("Sebagai") if "Sebagai" in headers else -1
        
        # Kolom penandatangan menggunakan struktur asli pengguna (Penanda Tangan 1, 2, 3)
        idx_t1_nama = headers.index("Penanda Tangan 1")
        idx_t1_jab = idx_t1_nama + 1
        
        idx_t2_nama = headers.index("Penanda Tangan 2")
        idx_t2_jab = idx_t2_nama + 1
        
        idx_t3_nama = headers.index("Penanda Tangan 3")
        idx_t3_jab = idx_t3_nama + 1
    except ValueError as e:
        print(f"Error: Kolom wajib Excel tidak ditemukan. Detail error: {e}")
        return

    # Tempat menyimpan metadata untuk website database.js
    certificate_data_list = {}
    
    # Inisialisasi PowerPoint COM secara Headless (Latar Belakang)
    powerpoint_app = None
    powerpoint_initialized = False

    # Loop setiap baris data sertifikat (Mulai dari Baris 3 karena Baris 1-2 adalah Header gabungan)
    row_count = ws.max_row
    
    # Simpan daftar QR temporary untuk dihapus nanti
    temp_files = []
    
    # Catat berapa banyak sertifikat baru yang benar-benar dibuat
    new_certs_processed = 0

    for row_idx in range(3, row_count + 1):
        id_cert = ws.cell(row=row_idx, column=idx_id + 1).value
        penerima = ws.cell(row=row_idx, column=idx_penerima + 1).value
        kegiatan = ws.cell(row=row_idx, column=idx_kegiatan + 1).value
        sebagai = ws.cell(row=row_idx, column=idx_sebagai + 1).value if idx_sebagai != -1 else "Peserta"
        timestamp = ws.cell(row=row_idx, column=idx_timestamp + 1).value
        dokumen_link = ws.cell(row=row_idx, column=idx_dokumen + 1).value

        # Baca data penandatangan dari kolom yang sesuai
        ttd1_nama = ws.cell(row=row_idx, column=idx_t1_nama + 1).value or ""
        ttd1_jab = ws.cell(row=row_idx, column=idx_t1_jab + 1).value or ""
        ttd2_nama = ws.cell(row=row_idx, column=idx_t2_nama + 1).value or ""
        ttd2_jab = ws.cell(row=row_idx, column=idx_t2_jab + 1).value or ""
        ttd3_nama = ws.cell(row=row_idx, column=idx_t3_nama + 1).value or ""
        ttd3_jab = ws.cell(row=row_idx, column=idx_t3_jab + 1).value or ""

        # Pastikan data utama tidak kosong
        if not id_cert or not penerima:
            continue
            
        # Tentukan path file PDF output
        safe_id = sanitize_filename(str(id_cert))
        pdf_filename = f"{safe_id}.pdf"
        output_pdf_path = os.path.abspath(os.path.join(PDF_DIR, pdf_filename))
        pdf_url = f"{GITHUB_BASE_URL}/pdf/{pdf_filename}"
        
        # Bersihkan format string tanggal
        if hasattr(timestamp, 'strftime'):
            date_str = timestamp.strftime('%d %B %Y')
        else:
            date_str = str(timestamp) if timestamp else "Juli 2026"

        # Susun daftar penandatangan untuk database website
        signers_list = []
        if ttd1_nama:
            signers_list.append({"name": str(ttd1_nama), "role": str(ttd1_jab)})
        if ttd2_nama:
            signers_list.append({"name": str(ttd2_nama), "role": str(ttd2_jab)})
        if ttd3_nama:
            signers_list.append({"name": str(ttd3_nama), "role": str(ttd3_jab)})
            
        # Default jika kosong sama sekali
        if not signers_list:
            signers_list.append({
                "name": "Prof. Dr. dr. Siti Aminah, M.Kes",
                "role": "Dekan FIKES - UF"
            })

        # Simpan metadata untuk database.js (selalu diupdate agar database web lengkap)
        certificate_data_list[str(id_cert)] = {
            "name": str(penerima),
            "role": str(sebagai),
            "activity": str(kegiatan),
            "date": date_str,
            "pdfUrl": pdf_url,
            "signers": signers_list
        }

        # JIKA file PDF sudah ada DAN kolom Excel Dokumen Link sudah terisi, LEWATKAN
        if os.path.exists(output_pdf_path) and dokumen_link:
            continue
            
        print(f"Memproses sertifikat baru [{row_idx - 2}]: {id_cert} - {penerima}")
        new_certs_processed += 1

        # Inisialisasi PowerPoint hanya ketika ada sertifikat baru yang perlu diproses
        if not powerpoint_initialized:
            print("Membuka PowerPoint...")
            try:
                powerpoint_app = win32com.client.Dispatch("PowerPoint.Application")
                try:
                    powerpoint_app.Visible = False
                except Exception:
                    powerpoint_app.Visible = True
                powerpoint_initialized = True
            except Exception as e:
                print(f"Gagal membuka PowerPoint via COM. Pastikan MS PowerPoint terinstall. Error: {e}")
                return

        # 1. Buat QR Code
        verification_url = f"{GITHUB_BASE_URL}/?id={id_cert}"
        temp_qr_path = os.path.join(BASE_DIR, f"temp_qr_{safe_id}.png")
        generate_qr_code(verification_url, temp_qr_path)
        temp_files.append(temp_qr_path)

        # 2. Proses Modifikasi PPTX menggunakan python-pptx
        prs = Presentation(PPTX_TEMPLATE_PATH)
        slide = prs.slides[0]
        
        replacements = {
            "{NomorSertifikat}": str(id_cert),
            "{Penerima}": str(penerima),
            "{Sebagai}": str(sebagai),
            "{Kegiatan}": str(kegiatan),
            "{Ttd1Nama}": str(ttd1_nama),
            "{Ttd1Jabatan}": str(ttd1_jab),
            "{Ttd2Nama}": str(ttd2_nama),
            "{Ttd2Jabatan}": str(ttd2_jab),
            "{Ttd3Nama}": str(ttd3_nama),
            "{Ttd3Jabatan}": str(ttd3_jab)
        }
        
        # Ganti teks placeholder dan cari posisi semua shape QR
        qr_placeholders = ["{QR}", "{QR1}", "{QR2}", "{QR3}"]
        qr_shape_boxes = {}
        
        for shape in list(slide.shapes):
            if shape.has_text_frame:
                replace_text_in_shape(shape, replacements)
                
                # Cek apakah shape berisi salah satu placeholder QR
                for qr_ph in qr_placeholders:
                    if qr_ph in shape.text:
                        qr_shape_boxes[qr_ph] = {
                            "left": shape.left,
                            "top": shape.top,
                            "width": shape.width,
                            "height": shape.height
                        }
                        # Hapus shape placeholder teks QR
                        sp = shape._element
                        sp.getparent().remove(sp)
                        break

        # Tempel gambar QR Code verifikasi di semua lokasi placeholder QR yang ditemukan
        for qr_ph, box in qr_shape_boxes.items():
            slide.shapes.add_picture(
                temp_qr_path, 
                box["left"], 
                box["top"], 
                box["width"], 
                box["height"]
            )

        # Simpan presentasi sementara
        temp_pptx_name = f"temp_merge_{safe_id}.pptx"
        temp_pptx_path = os.path.join(BASE_DIR, temp_pptx_name)
        prs.save(temp_pptx_path)
        temp_files.append(temp_pptx_path)

        # 3. Konversi PPTX ke PDF menggunakan PowerPoint COM
        try:
            presentation_obj = powerpoint_app.Presentations.Open(os.path.abspath(temp_pptx_path), WithWindow=False)
            presentation_obj.SaveAs(output_pdf_path, 32)
            presentation_obj.Close()
            
            # Tulis URL dokumen ke Excel
            ws.cell(row=row_idx, column=idx_dokumen + 1, value=pdf_url)
            print(f"-> Sukses membuat PDF: {pdf_filename}")
        except Exception as e:
            print(f"-> Gagal konversi PDF untuk {id_cert}. Error: {e}")

    # Tutup PowerPoint jika diinisialisasi
    if powerpoint_initialized and powerpoint_app:
        print("Menutup PowerPoint...")
        try:
            powerpoint_app.Quit()
        except:
            pass

    # Bersihkan file-file sementara
    if temp_files:
        print("Membersihkan file temporary...")
        for temp_file in temp_files:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
            except Exception as e:
                print(f"Gagal menghapus file temp {temp_file}: {e}")

    # Simpan kembali Excel jika ada perubahan
    if new_certs_processed > 0:
        try:
            wb.save(EXCEL_PATH)
            print("Excel database berhasil diupdate dengan link dokumen.")
        except PermissionError:
            print("\n[ERROR] GAGAL MENYIMPAN EXCEL TERBARU!")
            print("Data telah diproses namun link dokumen tidak bisa disimpan di Excel.")
            print("Pastikan Anda menutup file 'sertifikat.xlsx' sebelum menjalankan skrip.")
    else:
        print("Tidak ada sertifikat baru yang perlu dibuat.")

    # 4. Generate database.js untuk website verifikasi (selalu digenerate agar sinkron)
    database_js_path = os.path.join(OUTPUT_DIR, "database.js")
    try:
        with open(database_js_path, "w", encoding="utf-8") as f:
            f.write("// Database Sertifikat Terverifikasi\n")
            f.write("// File ini digenerate otomatis oleh mail_merge.py. Jangan edit manual!\n\n")
            f.write("const certificateDatabase = ")
            
            import json
            json.dump(certificate_data_list, f, indent=2, ensure_ascii=False)
            f.write(";\n")
        print(f"File database website berhasil diperbarui di: {database_js_path}")
    except Exception as e:
        print(f"Gagal menulis file database.js: {e}")

    print("\n=== Proses Selesai! ===")
    print(f"Total sertifikat baru diproses: {new_certs_processed}")
    print(f"Semua file siap diunggah berada di: {OUTPUT_DIR}")

if __name__ == "__main__":
    process_mail_merge()
